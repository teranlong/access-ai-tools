"""
pptx_to_md.py — Convert PowerPoint presentations to per-slide Markdown.

Extracts titles, text frames (preserving indent level), tables, and speaker notes.
"""

from __future__ import annotations

from pathlib import Path

from .base import ConversionResult, make_output_path


def convert(
    source: Path,
    output_root: Path,
    project_root: Path,
    **kwargs,
) -> ConversionResult:
    """Convert a PowerPoint file to structured Markdown."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError as exc:
        return ConversionResult(
            success=False,
            action="error",
            error=f"Missing dependency: {exc}. Run: pip install python-pptx",
        )

    out_path = make_output_path(source, output_root, project_root, ".md")
    sections: list[str] = [f"# Presentation: {source.name}\n"]

    try:
        prs = Presentation(str(source))

        for slide_num, slide in enumerate(prs.slides, start=1):
            title_text = ""
            body_lines: list[str] = []

            for shape in slide.shapes:
                # Title placeholder
                if shape.has_text_frame and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    continue

                if shape.has_text_frame:
                    tf = shape.text_frame
                    is_title = (
                        hasattr(shape, "placeholder_format")
                        and shape.placeholder_format is not None
                        and shape.placeholder_format.idx == 0
                    )
                    if is_title:
                        title_text = tf.text.strip()
                        continue

                    for para in tf.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        level = para.level or 0
                        indent = "  " * level
                        body_lines.append(f"{indent}- {text}")

                elif shape.has_table:
                    table = shape.table
                    rows_md: list[str] = []
                    for r_idx, row in enumerate(table.rows):
                        cells = [cell.text.replace("|", "\\|").replace("\n", " ").strip()
                                 for cell in row.cells]
                        rows_md.append("| " + " | ".join(cells) + " |")
                        if r_idx == 0:
                            rows_md.append("| " + " | ".join(["---"] * len(cells)) + " |")
                    body_lines.extend(rows_md)

            # Speaker notes
            notes_text = ""
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                notes_text = notes_tf.text.strip() if notes_tf else ""

            heading = f"## Slide {slide_num}: {title_text}" if title_text else f"## Slide {slide_num}"
            slide_section = heading + "\n"
            if body_lines:
                slide_section += "\n" + "\n".join(body_lines) + "\n"
            if notes_text:
                slide_section += f"\n**Speaker notes:** {notes_text}\n"

            sections.append(slide_section)

    except Exception as exc:
        out_path.write_text(
            f"# Conversion Error: {source.name}\n\n**Error:** `{exc}`\n",
            encoding="utf-8",
        )
        return ConversionResult(
            success=False,
            output_files=[out_path],
            action="error",
            error=str(exc),
        )

    out_path.write_text("\n---\n\n".join(sections), encoding="utf-8")

    return ConversionResult(
        success=True,
        output_files=[out_path],
        action="converted",
    )
